
import os
import logging
from typing import List, Dict, Optional
from datetime import datetime
from src.utils import setup_logging

logger = setup_logging(__name__)

# Check for COM availability (Windows only)
try:
    import win32com.client
    import pythoncom
    COM_AVAILABLE = True
except ImportError:
    COM_AVAILABLE = False

class PresentationGenerator:
    """Generates PowerPoint presentations"""
    
    def __init__(self, output_dir: str):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
        
    def generate(self, title: str, summary: str, references: List[Dict]) -> Optional[str]:
        """
        Generate a PPTX file.
        Uses COM if available for best results, otherwise creates a basic PPTX using python-pptx.
        """
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"presentation_{timestamp}.pptx"
        output_path = os.path.join(self.output_dir, filename)
        
        if COM_AVAILABLE:
            try:
                return self._generate_with_com(title, summary, references, output_path)
            except Exception as e:
                logger.error(f"COM generation failed, falling back to basic: {e}")
        
        return self._generate_basic(title, summary, references, output_path)

    def _generate_with_com(self, title: str, summary: str, references: List[Dict], path: str) -> str:
        """Windows-specific COM automation"""
        pythoncom.CoInitialize()
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True
        
        try:
            pres = ppt_app.Presentations.Add()
            
            # Slide 1: Title
            slide1 = pres.Slides.Add(1, 1) # ppLayoutTitle
            slide1.Shapes.Title.TextFrame.TextRange.Text = title
            
            # Slide 2: Summary
            slide2 = pres.Slides.Add(2, 2) # ppLayoutText
            slide2.Shapes.Title.TextFrame.TextRange.Text = "Summary"
            slide2.Shapes.Placeholders(2).TextFrame.TextRange.Text = summary
            
            # Simple placeholder for reference copying logic
            # In a full implementation, this opens source PPTs and copies slides.
            # Simplified here for the portfolio version.
            
            pres.SaveAs(os.path.abspath(path))
            return path
        except Exception as e:
            raise e
        finally:
            # Keep open or close depending on preference? usually close for automation, but open for demo.
            # pres.Close() 
            # ppt_app.Quit() # Don't quit if we want user to see it
            pass

    def _generate_basic(self, title: str, summary: str, references: List[Dict], path: str) -> str:
        """Cross-platform basic generation"""
        from pptx import Presentation
        
        prs = Presentation()
        
        # Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        # Content Slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Summary"
        slide.placeholders[1].text = summary
        
        prs.save(path)
        logger.info(f"Generated basic presentation at {path}")
        return path
