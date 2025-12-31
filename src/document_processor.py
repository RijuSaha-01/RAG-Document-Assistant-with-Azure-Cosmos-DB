
import os
import hashlib
from typing import List, Dict, Any, Generator
from datetime import datetime
from langchain.text_splitter import RecursiveCharacterTextSplitter
from src.utils import setup_logging

logger = setup_logging(__name__)

class DocumentProcessor:
    """Handles loading and chunking of documents (PDF, DOCX, PPTX)"""
    
    def __init__(self, chunk_size=800, chunk_overlap=100):
        # We use a recursive splitter to preserve paragraph and sentence boundaries
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""]
        )

    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process a single file and return chunks with metadata.
        """
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        
        text_content = ""
        metadata = {"source": filename, "type": ext}
        
        try:
            # Format-specific extraction
            if ext == '.pdf':
                text_content = self._read_pdf(file_path)
            elif ext == '.docx':
                text_content = self._read_docx(file_path)
            elif ext in ['.pptx', '.pptm']:
                text_content = self._read_pptx(file_path)
            else:
                logger.warning(f"Unsupported file type: {ext}")
                return []
                
            if not text_content:
                logger.warning(f"No text extracted from {filename}")
                return []
                
            # Create chunks using LangChain
            chunks = self.text_splitter.create_documents([text_content])
            
            # Format for Vector Store ingestion
            processed_chunks = []
            for i, chunk in enumerate(chunks):
                # Generate a unique ID for each chunk for upserting
                doc_id = hashlib.md5(f"{filename}_{i}".encode()).hexdigest()
                processed_chunks.append({
                    "id": doc_id,
                    "text": chunk.page_content,
                    "embedding": [], # Populated by the pipeline
                    "metadata": {
                        **metadata,
                        "chunk_index": i
                    }
                })
                
            logger.info(f"Processed {filename}: {len(processed_chunks)} chunks")
            return processed_chunks
            
        except Exception as e:
            logger.error(f"Error processing file {filename}: {e}")
            return []

    def _read_pdf(self, path: str) -> str:
        """Read PDF using pypdf"""
        import pypdf
        text = []
        try:
            reader = pypdf.PdfReader(path)
            for page in reader.pages:
                text.append(page.extract_text() or "")
            return "\n".join(text)
        except Exception as e:
            logger.error(f"PDF read error: {e}")
            return ""

    def _read_docx(self, path: str) -> str:
        """Read DOCX using python-docx"""
        import docx
        try:
            doc = docx.Document(path)
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            logger.error(f"DOCX read error: {e}")
            return ""

    def _read_pptx(self, path: str) -> str:
        """Read PPTX using python-pptx"""
        from pptx import Presentation
        try:
            prs = Presentation(path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            logger.error(f"PPTX read error: {e}")
            return ""
