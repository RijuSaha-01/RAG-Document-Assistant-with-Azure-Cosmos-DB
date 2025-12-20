"""
Enhanced Document Processor for PDF, DOCX, and PPTX files
Supports text extraction, metadata enrichment, and chunking
"""

import os
import logging
from typing import List, Dict, Any, Generator, Optional
import hashlib
from datetime import datetime

# PDF processing
try:
    import fitz  # PyMuPDF
    PDF_LIBRARY = 'pymupdf'
except ImportError:
    try:
        import pypdf
        PDF_LIBRARY = 'pypdf'
    except ImportError:
        PDF_LIBRARY = None

# Document processing
from pptx import Presentation
from docx import Document
from PIL import Image
import pytesseract

# Text processing
from langchain.text_splitter import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Enhanced document processor with support for multiple formats"""
    
    def __init__(self, chunk_size: int = 800, chunk_overlap: int = 150):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n\n", "\n\n", "\n", ". ", " ", ""]
        )
        
    def process_file(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        """Process a file and yield document chunks with metadata"""
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return
            
        filename = os.path.basename(file_path)
        ext = os.path.splitext(file_path)[1].lower()
        
        logger.info(f"Processing file: {filename} (type: {ext})")
        
        try:
            if ext == ".pdf":
                yield from self._process_pdf(file_path)
            elif ext in [".pptx", ".pptm"]:
                yield from self._process_pptx(file_path)
            elif ext in [".docx", ".docm"]:
                yield from self._process_docx(file_path)
            else:
                logger.warning(f"Unsupported file type: {ext}")
                
        except Exception as e:
            logger.error(f"Error processing {filename}: {e}")
    
    def _process_pdf(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        """Process PDF file and extract text with page information"""
        filename = os.path.basename(file_path)
        
        if PDF_LIBRARY == 'pymupdf':
            yield from self._process_pdf_pymupdf(file_path)
        elif PDF_LIBRARY == 'pypdf':
            yield from self._process_pdf_pypdf(file_path)
        else:
            logger.error("No PDF library available. Install PyMuPDF or pypdf.")
            return
    
    def _process_pdf_pymupdf(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        """Process PDF using PyMuPDF (fitz)"""
        filename = os.path.basename(file_path)
        
        try:
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                
                if text.strip():
                    # Create chunks from page text
                    chunks = self.text_splitter.split_text(text)
                    
                    for i, chunk in enumerate(chunks):
                        if chunk.strip():
                            yield {
                                'text': chunk.strip(),
                                'metadata': {
                                    'source': filename,
                                    'page': page_num + 1,
                                    'chunk_index': i,
                                    'file_type': 'pdf',
                                    'total_pages': len(doc),
                                    'processed_at': datetime.utcnow().isoformat()
                                }
                            }
            
            doc.close()
            logger.info(f"Successfully processed PDF: {filename} ({len(doc)} pages)")
            
        except Exception as e:
            logger.error(f"Error processing PDF with PyMuPDF: {e}")
    
    def _process_pdf_pypdf(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        """Process PDF using pypdf"""
        filename = os.path.basename(file_path)
        
        try:
            import pypdf
            
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    
                    if text.strip():
                        # Create chunks from page text
                        chunks = self.text_splitter.split_text(text)
                        
                        for i, chunk in enumerate(chunks):
                            if chunk.strip():
                                yield {
                                    'text': chunk.strip(),
                                    'metadata': {
                                        'source': filename,
                                        'page': page_num + 1,
                                        'chunk_index': i,
                                        'file_type': 'pdf',
                                        'total_pages': len(pdf_reader.pages),
                                        'processed_at': datetime.utcnow().isoformat()
                                    }
                                }
                
                logger.info(f"Successfully processed PDF: {filename} ({len(pdf_reader.pages)} pages)")
                
        except Exception as e:
            logger.error(f"Error processing PDF with pypdf: {e}")
    
    def _process_pptx(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        """Process PowerPoint file and extract text with slide information"""
        filename = os.path.basename(file_path)
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    full_text = "\n".join(slide_text)
                    
                    # Create chunks from slide text
                    chunks = self.text_splitter.split_text(full_text)
                    
                    for i, chunk in enumerate(chunks):
                        if chunk.strip():
                            yield {
                                'text': chunk.strip(),
                                'metadata': {
                                    'source': filename,
                                    'slide': slide_num + 1,
                                    'chunk_index': i,
                                    'file_type': 'pptx',
                                    'total_slides': len(prs.slides),
                                    'processed_at': datetime.utcnow().isoformat()
                                }
                            }
            
            logger.info(f"Successfully processed PPTX: {filename} ({len(prs.slides)} slides)")
            
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
    
    def _process_docx(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        """Process Word document and extract text with paragraph information"""
        filename = os.path.basename(file_path)
        
        try:
            doc = Document(file_path)
            
            # Extract text from paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        paragraphs.append(" | ".join(row_text))
            
            if paragraphs:
                full_text = "\n\n".join(paragraphs)
                
                # Create chunks from document text
                chunks = self.text_splitter.split_text(full_text)
                
                for i, chunk in enumerate(chunks):
                    if chunk.strip():
                        yield {
                            'text': chunk.strip(),
                            'metadata': {
                                'source': filename,
                                'chunk_index': i,
                                'file_type': 'docx',
                                'total_paragraphs': len(paragraphs),
                                'processed_at': datetime.utcnow().isoformat()
                            }
                        }
            
            logger.info(f"Successfully processed DOCX: {filename} ({len(paragraphs)} paragraphs)")
            
        except Exception as e:
            logger.error(f"Error processing DOCX: {e}")
    
    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract basic metadata from file"""
        try:
            stat = os.stat(file_path)
            filename = os.path.basename(file_path)
            ext = os.path.splitext(file_path)[1].lower()
            
            metadata = {
                'filename': filename,
                'file_type': ext[1:] if ext else 'unknown',
                'file_size': stat.st_size,
                'created_at': datetime.fromtimestamp(stat.st_ctime).isoformat(),
                'modified_at': datetime.fromtimestamp(stat.st_mtime).isoformat(),
                'file_hash': self._calculate_file_hash(file_path)
            }
            
            # Add format-specific metadata
            if ext == '.pdf':
                metadata.update(self._get_pdf_metadata(file_path))
            elif ext in ['.pptx', '.pptm']:
                metadata.update(self._get_pptx_metadata(file_path))
            elif ext in ['.docx', '.docm']:
                metadata.update(self._get_docx_metadata(file_path))
            
            return metadata
            
        except Exception as e:
            logger.error(f"Error extracting metadata from {file_path}: {e}")
            return {'filename': os.path.basename(file_path), 'error': str(e)}
    
    def _calculate_file_hash(self, file_path: str) -> str:
        """Calculate MD5 hash of file for deduplication"""
        try:
            hash_md5 = hashlib.md5()
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception:
            return ""
    
    def _get_pdf_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract PDF-specific metadata"""
        metadata = {}
        try:
            if PDF_LIBRARY == 'pymupdf':
                doc = fitz.open(file_path)
                metadata.update({
                    'page_count': len(doc),
                    'pdf_metadata': doc.metadata
                })
                doc.close()
            elif PDF_LIBRARY == 'pypdf':
                import pypdf
                with open(file_path, 'rb') as file:
                    pdf_reader = pypdf.PdfReader(file)
                    metadata.update({
                        'page_count': len(pdf_reader.pages),
                        'pdf_metadata': dict(pdf_reader.metadata) if pdf_reader.metadata else {}
                    })
        except Exception as e:
            logger.warning(f"Could not extract PDF metadata: {e}")
        
        return metadata
    
    def _get_pptx_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract PPTX-specific metadata"""
        metadata = {}
        try:
            prs = Presentation(file_path)
            metadata.update({
                'slide_count': len(prs.slides),
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height
            })
            
            # Extract core properties if available
            if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
                metadata['title'] = prs.core_properties.title
            if hasattr(prs.core_properties, 'author') and prs.core_properties.author:
                metadata['author'] = prs.core_properties.author
            if hasattr(prs.core_properties, 'subject') and prs.core_properties.subject:
                metadata['subject'] = prs.core_properties.subject
                
        except Exception as e:
            logger.warning(f"Could not extract PPTX metadata: {e}")
        
        return metadata
    
    def _get_docx_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract DOCX-specific metadata"""
        metadata = {}
        try:
            doc = Document(file_path)
            
            # Count elements
            metadata.update({
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables)
            })
            
            # Extract core properties if available
            if hasattr(doc.core_properties, 'title') and doc.core_properties.title:
                metadata['title'] = doc.core_properties.title
            if hasattr(doc.core_properties, 'author') and doc.core_properties.author:
                metadata['author'] = doc.core_properties.author
            if hasattr(doc.core_properties, 'subject') and doc.core_properties.subject:
                metadata['subject'] = doc.core_properties.subject
                
        except Exception as e:
            logger.warning(f"Could not extract DOCX metadata: {e}")
        
        return metadata