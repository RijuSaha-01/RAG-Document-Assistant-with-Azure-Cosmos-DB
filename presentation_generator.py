"""
PowerPoint Presentation Generator
Creates presentations from conversation data using COM integration when available
"""

import os
import logging
from typing import List, Dict, Any, Optional
from datetime import datetime
import tempfile

# PowerPoint creation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Windows COM Integration (optional)
try:
    import win32com.client
    import pythoncom
    COM_AVAILABLE = True
except ImportError:
    COM_AVAILABLE = False

logger = logging.getLogger(__name__)

class PresentationGenerator:
    """Generate PowerPoint presentations from conversation data"""
    
    def __init__(self, data_directory: str):
        self.data_directory = data_directory
        self.presentations_dir = os.path.join(data_directory, "generated_presentations")
        os.makedirs(self.presentations_dir, exist_ok=True)
        
        # COM integration
        self.com_available = COM_AVAILABLE
        self.ppt_app = None
        
        logger.info(f"PresentationGenerator initialized (COM available: {COM_AVAILABLE})")
    
    def create_presentation(self, query: str, response: str, source_references: List[Dict]) -> Optional[str]:
        """Create PowerPoint presentation from conversation data"""
        try:
            # Generate filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{timestamp}.pptx"
            output_path = os.path.join(self.presentations_dir, filename)
            
            # Try COM-based generation first (for slide copying)
            if self.com_available and source_references:
                com_path = self._create_com_presentation(query, response, source_references, output_path)
                if com_path:
                    return com_path
            
            # Fallback to python-pptx
            return self._create_basic_presentation(query, response, source_references, output_path)
            
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            return None
    
    def _create_com_presentation(self, query: str, response: str, source_references: List[Dict], output_path: str) -> Optional[str]:
        """Create presentation using COM with actual slide copying"""
        try:
            # Initialize COM
            if not self._initialize_com():
                return None
            
            # Create new presentation
            presentation = self.ppt_app.Presentations.Add()
            
            # Add title slide
            self._add_title_slide_com(presentation, query)
            
            # Add summary slide
            self._add_summary_slide_com(presentation, response)
            
            # Copy referenced slides efficiently
            copied_count = self._copy_slides_efficiently_com(presentation, source_references)
            
            # Save presentation
            abs_output_path = os.path.abspath(output_path)
            presentation.SaveAs(abs_output_path)
            presentation.Close()
            
            logger.info(f"Created COM presentation with {copied_count} copied slides: {output_path}")
            return abs_output_path
            
        except Exception as e:
            logger.error(f"Error creating COM presentation: {e}")
            return None
        finally:
            self._cleanup_com()
    
    def _initialize_com(self) -> bool:
        """Initialize PowerPoint COM application"""
        try:
            pythoncom.CoInitialize()
            self.ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            self.ppt_app.Visible = True
            return True
        except Exception as e:
            logger.error(f"Failed to initialize COM: {e}")
            return False
    
    def _cleanup_com(self):
        """Cleanup COM resources"""
        try:
            if self.ppt_app:
                self.ppt_app.Quit()
                self.ppt_app = None
            pythoncom.CoUninitialize()
        except Exception as e:
            logger.warning(f"Error during COM cleanup: {e}")
    
    def _add_title_slide_com(self, presentation, query: str):
        """Add title slide using COM"""
        try:
            slide = presentation.Slides.Add(1, 1)  # ppLayoutTitle
            slide.Shapes.Title.TextFrame.TextRange.Text = "Generated Presentation"
            
            if len(slide.Shapes.Placeholders) > 1:
                slide.Shapes.Placeholders(2).TextFrame.TextRange.Text = (
                    f"Query: {query}\n"
                    f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
                )
        except Exception as e:
            logger.error(f"Error adding title slide: {e}")
    
    def _add_summary_slide_com(self, presentation, response: str):
        """Add summary slide using COM"""
        try:
            slide = presentation.Slides.Add(2, 2)  # ppLayoutText
            slide.Shapes.Title.TextFrame.TextRange.Text = "Summary"
            
            # Truncate response if too long
            summary_text = response[:1000] + "..." if len(response) > 1000 else response
            slide.Shapes.Placeholders(2).TextFrame.TextRange.Text = summary_text
            
        except Exception as e:
            logger.error(f"Error adding summary slide: {e}")
    
    def _copy_slides_efficiently_com(self, presentation, source_references: List[Dict]) -> int:
        """Copy referenced slides efficiently by grouping by source file"""
        copied_count = 0
        
        # Group references by source file
        refs_by_source = {}
        for ref in source_references:
            source_name = ref['source_name']
            if source_name not in refs_by_source:
                refs_by_source[source_name] = []
            refs_by_source[source_name].append(ref)
        
        # Process each source file once
        for source_name, refs in refs_by_source.items():
            try:
                source_path = self._find_source_file(source_name)
                if not source_path:
                    logger.warning(f"Source file not found: {source_name}")
                    continue
                
                # Open source presentation once
                source_ppt = self.ppt_app.Presentations.Open(
                    os.path.abspath(source_path),
                    ReadOnly=True,
                    Untitled=False,
                    WithWindow=False
                )
                
                try:
                    # Copy all required slides from this source
                    for ref in refs:
                        slide_num = int(ref['page_or_slide'])
                        
                        # Check if slide exists
                        if slide_num > source_ppt.Slides.Count:
                            logger.warning(f"Slide {slide_num} not found in {source_name}")
                            continue
                        
                        # Copy slide
                        source_slide = source_ppt.Slides(slide_num)
                        source_slide.Copy()
                        
                        # Paste to target presentation
                        slide_index = presentation.Slides.Count + 1
                        presentation.Slides.Paste(slide_index)
                        
                        # Add source annotation
                        try:
                            pasted_slide = presentation.Slides(slide_index)
                            header_shape = pasted_slide.Shapes.AddTextbox(1, 20, 10, 680, 30)
                            header_shape.TextFrame.TextRange.Text = f"Source: {source_name} - Slide {slide_num}"
                            header_shape.TextFrame.TextRange.Font.Size = 12
                            header_shape.TextFrame.TextRange.Font.Bold = True
                        except Exception:
                            pass  # Ignore header errors
                        
                        copied_count += 1
                        logger.info(f"Copied slide {slide_num} from {source_name}")
                        
                finally:
                    source_ppt.Close()
                    
            except Exception as e:
                logger.error(f"Error copying slides from {source_name}: {e}")
                continue
        
        return copied_count

    def _copy_referenced_slide_com(self, presentation, reference: Dict) -> bool:
        """Copy a referenced slide from source file"""
        try:
            source_name = reference['source_name']
            slide_num = int(reference['page_or_slide'])
            
            # Find source file in data directory
            source_path = self._find_source_file(source_name)
            if not source_path:
                logger.warning(f"Source file not found: {source_name}")
                return False
            
            # Open source presentation
            source_ppt = self.ppt_app.Presentations.Open(
                os.path.abspath(source_path),
                ReadOnly=True,
                Untitled=False,
                WithWindow=False
            )
            
            try:
                # Check if slide exists
                if slide_num > source_ppt.Slides.Count:
                    logger.warning(f"Slide {slide_num} not found in {source_name}")
                    return False
                
                # Copy slide
                source_slide = source_ppt.Slides(slide_num)
                source_slide.Copy()
                
                # Paste to target presentation
                slide_index = presentation.Slides.Count + 1
                presentation.Slides.Paste(slide_index)
                
                # Add source annotation
                pasted_slide = presentation.Slides(slide_index)
                try:
                    header_shape = pasted_slide.Shapes.AddTextbox(1, 20, 10, 680, 30)
                    header_shape.TextFrame.TextRange.Text = f"Source: {source_name} - Slide {slide_num}"
                    header_shape.TextFrame.TextRange.Font.Size = 12
                    header_shape.TextFrame.TextRange.Font.Bold = True
                except Exception:
                    pass  # Ignore header errors
                
                logger.info(f"Copied slide {slide_num} from {source_name}")
                return True
                
            finally:
                source_ppt.Close()
                
        except Exception as e:
            logger.error(f"Error copying slide from {reference}: {e}")
            return False
    
    def _find_source_file(self, filename: str) -> Optional[str]:
        """Find source file in data directory"""
        # First try exact match
        exact_path = os.path.join(self.data_directory, filename)
        if os.path.exists(exact_path) and exact_path.lower().endswith(('.pptx', '.pptm')):
            return exact_path
        
        # Search in data directory with fuzzy matching
        for root, dirs, files in os.walk(self.data_directory):
            for file in files:
                if file_path := os.path.join(root, file):
                    if file_path.lower().endswith(('.pptx', '.pptm')):
                        # Try various matching strategies
                        if (file.lower() == filename.lower() or 
                            filename.lower() in file.lower() or
                            file.lower() in filename.lower() or
                            os.path.splitext(file)[0].lower() == os.path.splitext(filename)[0].lower()):
                            return file_path
        
        logger.warning(f"Could not find source file: {filename}")
        return None
    
    def _create_basic_presentation(self, query: str, response: str, source_references: List[Dict], output_path: str) -> str:
        """Create basic presentation using python-pptx"""
        try:
            # Create presentation
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "Generated Presentation"
            subtitle.text = f"Query: {query}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            
            # Summary slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Summary"
            
            # Split response into bullet points
            response_lines = response.split('\n')
            bullet_text = ""
            for line in response_lines[:10]:  # Limit to 10 lines
                if line.strip():
                    bullet_text += f"• {line.strip()}\n"
            
            content.text = bullet_text
            
            # Reference slides
            if source_references:
                slide = prs.slides.add_slide(bullet_slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = "Referenced Sources"
                
                ref_text = ""
                for ref in source_references:
                    ref_text += f"• {ref['source_name']} - {ref['reference_type'].title()} {ref['page_or_slide']}\n"
                
                content.text = ref_text
            
            # Save presentation
            prs.save(output_path)
            logger.info(f"Created basic presentation: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Error creating basic presentation: {e}")
            raise